import re
from collections import defaultdict
from pathlib import Path
from typing import Dict, Tuple

from pptx import Presentation
from pptx.util import Inches
from PIL import Image

# --- Configuration: SET YOUR PARAMETERS HERE ---
# 1. Define the two keywords that identify your image pairs.
KEYWORD_1: str = 'Basic_ParallelSummary'
KEYWORD_2: str = 'Basic_RasterPlot'

# 2. Set the TOP-LEVEL folder to start the recursive search.
#    The script will search this folder and ALL subfolders within it.
INPUT_FOLDER: str = r"C:\Users\maxyc\PycharmProjects\Ratatouille\cuisine\PassivePuff_JuxtaCellular_FromJS_202509"

# 3. Define slide and quality settings.
TARGET_WIDTH_INCHES: float = 13.333
TARGET_HEIGHT_INCHES: float = 7.5
DPI: int = 900
# ---------------------------------------------------

def find_image_pairs(input_dir: Path, kw1: str, kw2: str) -> Dict[str, Dict[str, Path]]:
    """
    Recursively scans a directory and its subdirectories to group images
    by UID based on the provided keywords.
    """
    image_pairs = defaultdict(dict)
    pattern = re.compile(rf'({kw1}|{kw2})_([a-zA-Z0-9_-]+)\.png')
    
    print(f"🖼️  Recursively scanning for '{kw1}' and '{kw2}' pairs in: {input_dir}")
    
    # CHANGE: Use .rglob('*.png') for a recursive search through all subfolders.
    for file_path in input_dir.rglob('*.png'):
        match = pattern.match(file_path.name)
        if match:
            fig_type, uid = match.groups()
            # Store the full path to the file.
            image_pairs[uid][fig_type] = file_path
            
    return image_pairs

def process_pair(uid: str, paths: Dict[str, Path], kw1: str, kw2: str, canvas_size: Tuple[int, int]):
    """
    Combines a pair of images, scaling them together to optimally fill the canvas
    while preserving individual aspect ratios and handling transparency correctly.
    """
    try:
        print(f"  - Processing pair for UID: {uid}")
        
        canvas_width, canvas_height = canvas_size
        
        img1 = Image.open(paths[kw1]).convert("RGBA")
        img2 = Image.open(paths[kw2]).convert("RGBA")

        ar1 = img1.width / img1.height
        ar2 = img2.width / img2.height
        final_height = min(canvas_height, canvas_width / (ar1 + ar2))
        final_height = int(final_height)
        new_width1 = int(final_height * ar1)
        new_width2 = int(final_height * ar2)
        img1_resized = img1.resize((new_width1, final_height), Image.Resampling.LANCZOS)
        img2_resized = img2.resize((new_width2, final_height), Image.Resampling.LANCZOS)

        canvas = Image.new('RGBA', canvas_size, (255, 255, 255, 255))

        total_content_width = new_width1 + new_width2
        start_x = (canvas_width - total_content_width) // 2
        start_y = (canvas_height - final_height) // 2
        
        canvas.paste(img1_resized, (start_x, start_y), img1_resized)
        canvas.paste(img2_resized, (start_x + new_width1, start_y), img2_resized)
        
        # The combined image is saved in the same folder as the first image of the pair.
        output_dir = paths[kw1].parent
        output_path = output_dir / f"combined_{uid}.png"
        canvas.save(output_path, dpi=(DPI, DPI))
        print(f"    ✅ Saved to {output_path}")

    except Exception as e:
        print(f"    ❌ Failed to process pair {uid}. Error: {e}")

def create_presentation(IMAGE_FOLDER: str, PPTX_FILENAME: str):
    """
    Finds all combined images and inserts them into a new PowerPoint presentation.
    """
    image_dir = Path(IMAGE_FOLDER)
    if not image_dir.is_dir():
        print(f"❌ Error: Image folder not found at '{IMAGE_FOLDER}'")
        return

    # 1. Recursively find all 'combined_*.png' images and sort them.
    print(f"🔎 Finding 'combined_*.png' files in '{image_dir}'...")
    image_paths = sorted(list(image_dir.rglob('combined_*.png')))

    if not image_paths:
        print("❌ No 'combined_*.png' images found. Please run the previous script first.")
        return

    print(f"Found {len(image_paths)} images to add.")

    # 2. Create a new 16:9 presentation.
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6] # Index 6 is typically the 'Blank' layout

    # 3. Loop through each image and add it to a new slide.
    for image_path in image_paths:
        print(f"  + Adding {image_path.name}")
        slide = prs.slides.add_slide(blank_slide_layout)

        # 4. Insert the picture, scaled to the slide width.
        #    The left and top are set to 0 initially.
        pic = slide.shapes.add_picture(
            str(image_path),
            left=Inches(0),
            top=Inches(0),
            width=prs.slide_width
        )

        # 5. Center the image vertically.
        #    The library automatically calculates the picture's height to keep the aspect ratio.
        #    We calculate the offset needed to center it on the slide.
        top_offset = (prs.slide_height - pic.height) / 2
        pic.top = int(top_offset)

    # 6. Save the final presentation.
    try:
        output_path = image_dir / PPTX_FILENAME
        prs.save(output_path)
        print(f"\n✅ Successfully created presentation: {output_path}")
    except Exception as e:
        print(f"\n❌ Failed to save presentation. Error: {e}")


def main():
    """Main function to run the script."""
    input_dir = Path(INPUT_FOLDER)
    if not input_dir.is_dir():
        print(f"Error: Input folder not found at '{INPUT_FOLDER}'")
        return
        
    canvas_size = (int(TARGET_WIDTH_INCHES * DPI), int(TARGET_HEIGHT_INCHES * DPI))
    image_pairs = find_image_pairs(input_dir, KEYWORD_1, KEYWORD_2)
    
    if not image_pairs:
        print("No matching image files found. Check your keywords and INPUT_FOLDER.")
        return

    complete_pairs = {
        uid: files for uid, files in image_pairs.items() 
        if KEYWORD_1 in files and KEYWORD_2 in files
    }
    
    if not complete_pairs:
        print(f"Found images, but no complete pairs of ('{KEYWORD_1}', '{KEYWORD_2}').")
        return

    print(f"Found {len(complete_pairs)} complete pairs. Starting combination process...")
    
    for uid, paths in complete_pairs.items():
        process_pair(uid, paths, KEYWORD_1, KEYWORD_2, canvas_size)
        
    print("\nProcessing complete!")

if __name__ == "__main__":
    main()
    create_presentation(r"C:\Users\maxyc\PycharmProjects\Ratatouille\cuisine\PassivePuff_JuxtaCellular_FromJS_202509\PYR_JUX", "PYR_JUX.pptx")
    create_presentation(r"C:\Users\maxyc\PycharmProjects\Ratatouille\cuisine\PassivePuff_JuxtaCellular_FromJS_202509\PV_JUX", "PV_JUX.pptx")
    create_presentation(r"C:\Users\maxyc\PycharmProjects\Ratatouille\cuisine\PassivePuff_JuxtaCellular_FromJS_202509\SST_JUX", "SST_JUX.pptx")
    create_presentation(r"C:\Users\maxyc\PycharmProjects\Ratatouille\cuisine\PassivePuff_JuxtaCellular_FromJS_202509\SST_WC", "SST_WC.pptx")